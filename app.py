import streamlit as st
import os
import io
import re
import json
import base64
import requests
from urllib.parse import urlparse

from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_groq import ChatGroq
from langchain.agents import create_agent
from tavily import TavilyClient

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# =====================================================================
# PAGE CONFIG + LIGHT UI/UX POLISH
# =====================================================================
st.set_page_config(page_title="AI PPT Generator", page_icon="🎨", layout="wide")

st.markdown(
    """
    <style>
        .block-container {padding-top: 2rem; padding-bottom: 3rem;}
        .stButton>button {
            border-radius: 10px; padding: 0.55rem 1.1rem; font-weight: 600;
            border: none; background: linear-gradient(90deg,#4f46e5,#7c3aed); color: white;
        }
        .stButton>button:hover {opacity: 0.9;}
        .stDownloadButton>button {border-radius: 10px; font-weight: 600; width: 100%;}
        .app-subtitle {color:#64748b; margin-top:-10px; margin-bottom:1.4rem;}
        .tpl-card {border-radius:12px; overflow:hidden; box-shadow:0 2px 10px rgba(0,0,0,0.08);}
    </style>
    """,
    unsafe_allow_html=True,
)

st.title("🎨 AI PPT Generator")
st.markdown('<p class="app-subtitle">Describe your topic, pick a template, get a full deck — preview as HTML, download as PPTX or HTML.</p>', unsafe_allow_html=True)

GOOGLE_API_KEY = st.sidebar.text_input("Google Api Key", type="password")
GROQ_API_KEY = st.sidebar.text_input("GROQ Api Key (optional)", type="password")
TAVILY_API_KEY = st.sidebar.text_input("TAVILY Api Key (optional)", type="password")
if not GOOGLE_API_KEY:
    st.warning("Provide Google API key")

# =====================================================================
# 5 TEMPLATE DESIGNS
# =====================================================================
TEMPLATES = {
    "minimal_light": {
        "name": "Minimal Light",
        "tagline": "Clean, spacious, editorial",
        "bg": "#FFFFFF", "text": "#1e293b", "accent": "#4f46e5", "secondary": "#f1f5f9",
        "font_title": "'Poppins', sans-serif", "font_body": "'Inter', sans-serif",
        "gfont": "Poppins:wght@600;700|Inter:wght@400;500",
    },
    "dark_professional": {
        "name": "Dark Professional",
        "tagline": "Bold, high-contrast, executive",
        "bg": "#0f172a", "text": "#f1f5f9", "accent": "#facc15", "secondary": "#1e293b",
        "font_title": "'Sora', sans-serif", "font_body": "'Inter', sans-serif",
        "gfont": "Sora:wght@600;700|Inter:wght@400;500",
    },
    "gradient_modern": {
        "name": "Gradient Modern",
        "tagline": "Vibrant, energetic, startup-y",
        "bg": "#FFFFFF", "text": "#ffffff", "accent": "#fbbf24", "secondary": "rgba(255,255,255,0.15)",
        "gradient": ("#6366f1", "#8b5cf6"),
        "font_title": "'Poppins', sans-serif", "font_body": "'Inter', sans-serif",
        "gfont": "Poppins:wght@600;700|Inter:wght@400;500",
    },
    "corporate_blue": {
        "name": "Corporate Blue",
        "tagline": "Trustworthy, structured, formal",
        "bg": "#FFFFFF", "text": "#0f172a", "accent": "#2563eb", "secondary": "#eff6ff",
        "font_title": "'Source Sans 3', sans-serif", "font_body": "'Source Sans 3', sans-serif",
        "gfont": "Source+Sans+3:wght@600;700;400",
    },
    "creative_pastel": {
        "name": "Creative Pastel",
        "tagline": "Playful, soft, approachable",
        "bg": "#fff7f0", "text": "#4a3728", "accent": "#f97316", "secondary": "#ffe8d6",
        "font_title": "'Fredoka', sans-serif", "font_body": "'Nunito', sans-serif",
        "gfont": "Fredoka:wght@600;700|Nunito:wght@400;600",
    },
}

if "template_key" not in st.session_state:
    st.session_state.template_key = "minimal_light"
if "slides" not in st.session_state:
    st.session_state.slides = None
if "slide_images" not in st.session_state:
    st.session_state.slide_images = {}
if "pptx_bytes" not in st.session_state:
    st.session_state.pptx_bytes = None

st.write("### 🎨 Choose a template")
cols = st.columns(5)
for i, key in enumerate(TEMPLATES):
    t = TEMPLATES[key]
    is_selected = st.session_state.template_key == key
    preview_bg = f"linear-gradient(135deg,{t['gradient'][0]},{t['gradient'][1]})" if "gradient" in t else t["bg"]
    border = "#4f46e5" if is_selected else "#e2e8f0"
    with cols[i]:
        st.markdown(
            f"""
            <div class="tpl-card" style="border:2px solid {border};">
              <div style="height:64px;background:{preview_bg};display:flex;align-items:center;justify-content:center;">
                <div style="width:45%;height:8px;background:{t['accent']};border-radius:4px;"></div>
              </div>
              <div style="padding:8px;text-align:center;background:white;">
                <div style="font-size:13px;font-weight:700;color:#1e293b;">{t['name']}</div>
                <div style="font-size:11px;color:#94a3b8;">{t['tagline']}</div>
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        if st.button("Select" if not is_selected else "✓ Selected", key=f"sel_{key}", width="stretch"):
            st.session_state.template_key = key
            st.rerun()

theme = TEMPLATES[st.session_state.template_key]

st.divider()
topic_prompt = st.text_area("Describe your presentation", placeholder="e.g. A pitch deck for an AI-powered fitness app targeting Gen Z users")
slide_count = st.slider("Number of slides", 5, 15, 8)
use_research = st.checkbox("Enrich content with live web research (uses Tavily)", value=False, disabled=not TAVILY_API_KEY)

img_col1, img_col2 = st.columns([1, 1])
with img_col1:
    include_images = st.checkbox(
        "🖼️ Generate AI images for slides (Gemini 2.5 Flash Image, ~500/day free)",
        value=False, disabled=not GOOGLE_API_KEY,
    )
IMAGE_STYLES = {
    "Flat vector illustration": "flat vector illustration style, clean shapes, simple color palette, no text in image",
    "Photorealistic": "photorealistic photo, natural lighting, high detail, no text in image",
    "3D render": "3D rendered illustration, soft studio lighting, modern, no text in image",
    "Minimal line art": "minimal line-art illustration, single accent color, lots of white space, no text in image",
    "Watercolor": "soft watercolor illustration style, gentle colors, artistic, no text in image",
}
with img_col2:
    image_style = st.selectbox("Image style", list(IMAGE_STYLES.keys()), disabled=not include_images)


# =====================================================================
# TOOL: optional Tavily research (used only if key + checkbox provided)
# =====================================================================
def search_topic_facts(query):
    """Fetch a few recent facts/stats about the given topic using Tavily,
    to help ground the generated slide content."""
    client = TavilyClient(api_key=TAVILY_API_KEY)
    return client.search(query, search_depth="basic", max_results=5)


# =====================================================================
# Gemini 2.5 Flash Image ("Nano Banana") — uses the SAME Google API key
# already entered above, no separate key/service needed. Free tier is
# roughly 500 image requests/day (subject to Google's current limits).
# =====================================================================
GEMINI_IMAGE_MODEL = "gemini-2.5-flash-image"


def generate_image_bytes(prompt: str, api_key: str, style_suffix: str = "") -> bytes:
    """Calls the Gemini image-generation REST endpoint directly and returns
    raw image bytes. Raises on failure (caller should catch and warn)."""
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_IMAGE_MODEL}:generateContent?key={api_key}"
    full_prompt = f"{prompt}. Style: {style_suffix}".strip()
    payload = {"contents": [{"parts": [{"text": full_prompt}]}]}
    resp = requests.post(url, json=payload, timeout=60)
    resp.raise_for_status()
    data = resp.json()
    candidates = data.get("candidates", [])
    if not candidates:
        raise ValueError(f"No candidates returned: {data}")
    parts = candidates[0].get("content", {}).get("parts", [])
    for part in parts:
        inline = part.get("inlineData") or part.get("inline_data")
        if inline and inline.get("data"):
            return base64.b64decode(inline["data"])
    raise ValueError("Model did not return image data (it may have replied with text only).")


def get_model():
    return ChatGoogleGenerativeAI(model="gemini-3.5-flash-lite", google_api_key=GOOGLE_API_KEY)


def get_agent(model1):
    tools = [search_topic_facts] if TAVILY_API_KEY else []
    return create_agent(model=model1, tools=tools)


SLIDE_JSON_INSTRUCTIONS = """You are an expert presentation content writer.
Create a slide-by-slide outline for a presentation deck.

STRICT OUTPUT RULES:
- Output ONLY a valid JSON array. No prose, no markdown fences, no explanations.
- Each element is an object with exactly these fields:
  "layout": one of "title", "content", "two_column", "quote", "section", "closing", "image_focus"
  "title": short slide title (string)
  "subtitle": optional subtitle or empty string
  "bullets": array of short strings (3-5 items for "content"/"two_column", empty array [] for other layouts)
  "notes": one-sentence speaker note (string)
  "image_prompt": a short (under 15 words), purely descriptive, literal visual prompt for an
    AI image generator (e.g. "laptop on a desk with a rising bar chart on screen, office background"),
    or an empty string "" if this slide doesn't need an image. {image_instruction}
- The FIRST slide must have layout "title" (deck title + subtitle, bullets: []).
- The LAST slide must have layout "closing" (thank you / next steps / contact, bullets: []).
- If the deck has more than 6 slides, include at least one "section" divider slide.
- Use layout "image_focus" (title + supporting caption + a required image_prompt, bullets: [])
  for at most one or two slides where a single strong visual should dominate the slide.
- Produce exactly {slide_count} slide objects total.
- Keep bullet text concise (under 12 words each).

Topic: {topic}
{research_context}
"""


def extract_json_array(raw_text: str):
    cleaned = re.sub(r"```(?:json)?", "", raw_text).replace("```", "").strip()
    try:
        return json.loads(cleaned)
    except Exception:
        start, end = cleaned.find("["), cleaned.rfind("]")
        if start != -1 and end != -1 and end > start:
            return json.loads(cleaned[start:end + 1])
        raise


if st.button("✨ Generate Presentation"):
    if not GOOGLE_API_KEY:
        st.error("Please provide your Google API key in the sidebar first.")
    elif not topic_prompt.strip():
        st.error("Please describe what your presentation should be about.")
    else:
        with st.spinner("Researching and writing your deck..."):
            research_context = ""
            if use_research and TAVILY_API_KEY:
                try:
                    results = search_topic_facts(topic_prompt)
                    snippets = [r.get("content", "")[:200] for r in results.get("results", [])[:5]]
                    if snippets:
                        research_context = "Recent context to consider (use for inspiration, don't quote verbatim):\n- " + "\n- ".join(snippets)
                except Exception as e:
                    st.warning(f"Web research skipped ({e}). Continuing without it.")

            image_instruction = (
                "Add an image_prompt for roughly half the slides where a visual would help."
                if include_images else
                "Leave image_prompt as an empty string for every slide (images are disabled)."
            )
            prompt_text = SLIDE_JSON_INSTRUCTIONS.format(
                slide_count=slide_count, topic=topic_prompt, research_context=research_context,
                image_instruction=image_instruction,
            )

            model1 = get_model()
            agent = get_agent(model1)
            response = agent.invoke({"messages": [{"role": "user", "content": prompt_text}]})
            try:
                raw = response["messages"][-1].content[-1]["text"]
            except (TypeError, KeyError, IndexError):
                raw_content = response["messages"][-1].content
                raw = raw_content if isinstance(raw_content, str) else str(raw_content)

            try:
                slides = extract_json_array(raw)
                assert isinstance(slides, list) and len(slides) > 0
            except Exception as e:
                st.error(f"Couldn't parse the generated content as JSON, please try again. Details: {e}")
                slides = None

            slide_images = {}
            if slides and include_images and GOOGLE_API_KEY:
                with st.spinner("Generating images for your slides..."):
                    for i, s in enumerate(slides):
                        prompt = (s.get("image_prompt") or "").strip()
                        if not prompt:
                            continue
                        try:
                            slide_images[i] = generate_image_bytes(
                                prompt, GOOGLE_API_KEY, style_suffix=IMAGE_STYLES[image_style]
                            )
                        except Exception as e:
                            st.warning(f"Image generation failed for slide {i + 1} ('{prompt[:40]}...'): {e}")

            if slides:
                st.session_state.slides = slides
                st.session_state.slide_images = slide_images
                st.session_state.pptx_bytes = None
        if st.session_state.slides:
            n_imgs = len(st.session_state.get("slide_images") or {})
            msg = f"Generated {len(st.session_state.slides)} slides!"
            if n_imgs:
                msg += f" ({n_imgs} image{'s' if n_imgs != 1 else ''} generated)"
            st.success(msg)


# =====================================================================
# HTML DECK RENDERER
# =====================================================================
def render_html_deck(slides, theme, slide_images=None):
    slide_images = slide_images or {}
    bg_css = f"linear-gradient(135deg,{theme['gradient'][0]},{theme['gradient'][1]})" if "gradient" in theme else theme["bg"]
    gfont_url = "https://fonts.googleapis.com/css2?family=" + theme["gfont"] + "&display=swap"

    def img_tag(i):
        img_bytes = slide_images.get(i)
        if not img_bytes:
            return None
        b64 = base64.b64encode(img_bytes).decode()
        return f'<img class="slide-img" src="data:image/png;base64,{b64}" />'

    slide_html_blocks = []
    for i, s in enumerate(slides):
        layout = s.get("layout", "content")
        title = s.get("title", "")
        subtitle = s.get("subtitle", "")
        bullets = s.get("bullets", []) or []
        image_html = img_tag(i)

        if layout == "image_focus" and image_html:
            body = f"""
              <div class="layout-image-focus">
                <div class="img-pane full">{image_html}</div>
                <div class="caption-bar">
                  <h2>{title}</h2>
                  {f'<p class="subtitle">{subtitle}</p>' if subtitle else ''}
                </div>
              </div>"""
        elif layout in ("title", "closing"):
            text_block = f"""
                <h1>{title}</h1>
                {f'<p class="subtitle">{subtitle}</p>' if subtitle else ''}
                <div class="accent-line"></div>"""
            if image_html and layout == "title":
                body = f'<div class="layout-split"><div class="text-pane">{text_block}</div><div class="img-pane">{image_html}</div></div>'
            else:
                body = f'<div class="layout-title">{text_block}</div>'
        elif layout == "section":
            body = f"""
              <div class="layout-section">
                <div class="accent-line"></div>
                <h2>{title}</h2>
                {f'<p class="subtitle">{subtitle}</p>' if subtitle else ''}
              </div>"""
        elif layout == "quote":
            body = f"""
              <div class="layout-quote">
                <div class="quote-mark">&ldquo;</div>
                <p class="quote-text">{title}</p>
                {f'<p class="subtitle">— {subtitle}</p>' if subtitle else ''}
              </div>"""
        elif layout == "two_column":
            half = (len(bullets) + 1) // 2
            left_items = "".join(f"<li>{b}</li>" for b in bullets[:half])
            right_items = "".join(f"<li>{b}</li>" for b in bullets[half:])
            body = f"""
              <div class="layout-content">
                <h2>{title}</h2>
                {f'<p class="subtitle">{subtitle}</p>' if subtitle else ''}
                <div class="two-col">
                  <ul>{left_items}</ul>
                  <ul>{right_items}</ul>
                </div>
              </div>"""
        else:  # "content"
            items = "".join(f"<li>{b}</li>" for b in bullets)
            text_block = f"""
                <h2>{title}</h2>
                {f'<p class="subtitle">{subtitle}</p>' if subtitle else ''}
                <ul>{items}</ul>"""
            if image_html:
                body = f'<div class="layout-split"><div class="text-pane">{text_block}</div><div class="img-pane">{image_html}</div></div>'
            else:
                body = f'<div class="layout-content">{text_block}</div>'

        active = "active" if i == 0 else ""
        extra_cls = "no-pad" if (layout == "image_focus" and image_html) else ""
        slide_html_blocks.append(f'<section class="slide {active} {extra_cls}" data-index="{i}">{body}</section>')

    slides_markup = "\n".join(slide_html_blocks)

    return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<link rel="stylesheet" href="{gfont_url}">
<style>
  :root {{
    --bg: {theme['bg']}; --text: {theme['text']}; --accent: {theme['accent']}; --secondary: {theme['secondary']};
    --font-title: {theme['font_title']}; --font-body: {theme['font_body']};
  }}
  * {{ box-sizing: border-box; }}
  body {{ margin:0; font-family: var(--font-body); background:#111; }}
  .deck-wrap {{ max-width: 960px; margin: 0 auto; }}
  .deck {{ position:relative; width:100%; aspect-ratio: 16/9; background:{bg_css}; color: var(--text);
           border-radius: 14px; overflow:hidden; box-shadow: 0 10px 40px rgba(0,0,0,0.4); }}
  .slide {{ position:absolute; inset:0; display:none; padding: 6% 8%; flex-direction:column; justify-content:center; }}
  .slide.active {{ display:flex; }}
  h1 {{ font-family: var(--font-title); font-size: clamp(28px,5vw,52px); margin:0 0 12px 0; font-weight:700; }}
  h2 {{ font-family: var(--font-title); font-size: clamp(22px,3.5vw,34px); margin:0 0 14px 0; font-weight:700; }}
  .subtitle {{ font-size: clamp(14px,1.6vw,18px); opacity:0.85; margin:0 0 10px 0; }}
  .accent-line {{ width:70px; height:5px; background:var(--accent); border-radius:4px; margin-top:16px; }}
  .layout-title {{ text-align:left; }}
  .layout-section {{ align-items:flex-start; }}
  .layout-section h2 {{ font-size: clamp(26px,4vw,40px); }}
  .layout-quote {{ align-items:center; text-align:center; }}
  .quote-mark {{ font-size:60px; color:var(--accent); line-height:1; font-family: var(--font-title); }}
  .quote-text {{ font-size: clamp(18px,2.6vw,28px); font-style:italic; max-width:80%; }}
  ul {{ font-size: clamp(13px,1.7vw,19px); line-height:1.7; padding-left:22px; margin:0; }}
  li {{ margin-bottom:8px; }}
  li::marker {{ color: var(--accent); }}
  .two-col {{ display:flex; gap:40px; }}
  .two-col ul {{ flex:1; }}
  .layout-split {{ display:flex; align-items:center; gap:5%; width:100%; height:100%; }}
  .layout-split .text-pane {{ flex:1.1; min-width:0; }}
  .layout-split .img-pane {{ flex:0.9; min-width:0; height:70%; }}
  .slide-img {{ width:100%; height:100%; object-fit:cover; border-radius:12px; box-shadow:0 8px 24px rgba(0,0,0,0.25); }}
  .slide.no-pad {{ padding:0; }}
  .layout-image-focus {{ position:relative; width:100%; height:100%; }}
  .layout-image-focus .img-pane.full {{ width:100%; height:100%; }}
  .layout-image-focus .img-pane.full .slide-img {{ border-radius:0; }}
  .layout-image-focus .caption-bar {{ position:absolute; left:0; right:0; bottom:0; padding:5% 8%;
      background:linear-gradient(to top, rgba(0,0,0,0.75), rgba(0,0,0,0)); }}
  .layout-image-focus .caption-bar h2, .layout-image-focus .caption-bar .subtitle {{ color:#fff; margin:0; }}
  .nav-bar {{ display:flex; align-items:center; justify-content:space-between; margin-top:14px; color:#e2e8f0; font-family:sans-serif; font-size:13px; }}
  .nav-btn {{ background: var(--accent); color:#111; border:none; padding:8px 16px; border-radius:8px; cursor:pointer; font-weight:600; }}
  .nav-btn:disabled {{ opacity:0.35; cursor:default; }}
  .counter {{ opacity:0.8; }}
</style>
</head>
<body>
<div class="deck-wrap">
  <div class="deck" id="deck">
    {slides_markup}
  </div>
  <div class="nav-bar">
    <button class="nav-btn" id="prevBtn" onclick="changeSlide(-1)">&larr; Prev</button>
    <span class="counter"><span id="counter">1</span> / {len(slides)}</span>
    <button class="nav-btn" id="nextBtn" onclick="changeSlide(1)">Next &rarr;</button>
  </div>
</div>
<script>
  let idx = 0;
  const total = {len(slides)};
  function showSlide(i) {{
    document.querySelectorAll('.slide').forEach((el,n) => el.classList.toggle('active', n===i));
    document.getElementById('counter').innerText = i+1;
    document.getElementById('prevBtn').disabled = i===0;
    document.getElementById('nextBtn').disabled = i===total-1;
  }}
  function changeSlide(delta) {{
    idx = Math.min(total-1, Math.max(0, idx+delta));
    showSlide(idx);
  }}
  document.addEventListener('keydown', (e) => {{
    if (e.key === 'ArrowRight') changeSlide(1);
    if (e.key === 'ArrowLeft') changeSlide(-1);
  }});
  showSlide(0);
</script>
</body>
</html>"""


# =====================================================================
# PPTX RENDERER (python-pptx) — same slide data as the HTML deck above
# =====================================================================
def hexclr(h):
    return RGBColor.from_string(h.lstrip("#").upper())


def set_slide_background(slide, theme, variant="normal"):
    fill = slide.background.fill
    if "gradient" in theme:
        fill.gradient()
        stops = fill.gradient_stops
        stops[0].color.rgb = hexclr(theme["gradient"][0])
        stops[0].position = 0.0
        stops[1].color.rgb = hexclr(theme["gradient"][1])
        stops[1].position = 1.0
        fill.gradient_angle = 45
    elif variant == "section":
        fill.solid()
        fill.fore_color.rgb = hexclr(theme["accent"])
    else:
        fill.solid()
        fill.fore_color.rgb = hexclr(theme["bg"])


def add_text(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, size, color, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(10)
    return tb


def add_accent_bar(slide, theme, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hexclr(theme["accent"])
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_picture(slide, image_bytes, left, top, width, height):
    """Embeds an image stretched to fill the given box. python-pptx has no
    built-in CSS-style object-fit:cover, so very wide/tall source images may
    look slightly stretched; Gemini image output is close to square by
    default so this is rarely noticeable in practice."""
    return slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(left), Inches(top), Inches(width), Inches(height))


def build_pptx(slides, theme, slide_images=None):
    slide_images = slide_images or {}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    text_clr = hexclr(theme["text"])
    accent_clr = hexclr(theme["accent"])
    on_accent_clr = hexclr(theme["bg"]) if theme.get("bg") not in (None, "") else RGBColor(0xFF, 0xFF, 0xFF)
    body_font = "Calibri"
    title_font = "Calibri"

    for idx, s in enumerate(slides):
        layout = s.get("layout", "content")
        title = s.get("title", "")
        subtitle = s.get("subtitle", "")
        bullets = s.get("bullets", []) or []
        notes = s.get("notes", "")
        image_bytes = slide_images.get(idx)

        slide = prs.slides.add_slide(blank_layout)

        if layout == "image_focus" and image_bytes:
            add_picture(slide, image_bytes, 0, 0, 13.333, 7.5)
            caption_h = 1.9
            cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.5 - caption_h), Inches(13.333), Inches(caption_h))
            cap.fill.solid()
            cap.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            cap.line.fill.background()
            cap.shadow.inherit = False
            add_text(slide, 0.9, 7.5 - caption_h + 0.35, 11.5, 0.9, title, 28, RGBColor(0xFF, 0xFF, 0xFF), bold=True, font_name=title_font)
            if subtitle:
                add_text(slide, 0.9, 7.5 - caption_h + 1.05, 11.5, 0.6, subtitle, 15, RGBColor(0xFF, 0xFF, 0xFF), font_name=body_font)

        elif layout == "section":
            set_slide_background(slide, theme, variant="section")
            title_clr = RGBColor(0xFF, 0xFF, 0xFF) if not theme.get("gradient") else text_clr
            add_accent_bar(slide, theme, 0.9, 3.0, 0.6, 0.12)
            add_text(slide, 0.9, 3.2, 11.5, 1.2, title, 34, title_clr, bold=True, font_name=title_font)
            if subtitle:
                add_text(slide, 0.9, 4.1, 11.5, 0.8, subtitle, 16, title_clr, font_name=body_font)

        elif layout in ("title", "closing"):
            set_slide_background(slide, theme)
            if image_bytes and layout == "title":
                add_text(slide, 0.9, 2.5, 5.6, 1.6, title, 34, text_clr, bold=True, font_name=title_font)
                if subtitle:
                    add_text(slide, 0.9, 3.7, 5.6, 0.8, subtitle, 16, text_clr, font_name=body_font)
                add_accent_bar(slide, theme, 0.9, 4.3, 0.9, 0.09)
                add_picture(slide, image_bytes, 7.0, 1.2, 5.4, 5.1)
            else:
                add_text(slide, 0.9, 2.7, 11.5, 1.6, title, 40, text_clr, bold=True, font_name=title_font)
                if subtitle:
                    add_text(slide, 0.9, 3.9, 11.5, 0.8, subtitle, 18, text_clr, font_name=body_font)
                add_accent_bar(slide, theme, 0.9, 4.5, 0.9, 0.09)

        elif layout == "quote":
            set_slide_background(slide, theme)
            add_text(slide, 1.2, 1.6, 10.9, 1.2, "\u201C", 70, accent_clr, bold=True, align=PP_ALIGN.CENTER, font_name=title_font)
            add_text(slide, 1.2, 2.8, 10.9, 1.8, title, 26, text_clr, italic=True, align=PP_ALIGN.CENTER, font_name=body_font)
            if subtitle:
                add_text(slide, 1.2, 4.6, 10.9, 0.6, f"— {subtitle}", 16, text_clr, align=PP_ALIGN.CENTER, font_name=body_font)

        elif layout == "two_column":
            set_slide_background(slide, theme)
            add_accent_bar(slide, theme, 0.0, 0.0, 0.18, 7.5)
            add_text(slide, 0.9, 0.6, 11.5, 1.0, title, 30, text_clr, bold=True, font_name=title_font)
            if subtitle:
                add_text(slide, 0.9, 1.35, 11.5, 0.6, subtitle, 15, text_clr, font_name=body_font)
            half = (len(bullets) + 1) // 2
            add_bullets(slide, 0.9, 2.2, 5.6, 4.6, bullets[:half], 17, text_clr, font_name=body_font)
            add_bullets(slide, 6.8, 2.2, 5.6, 4.6, bullets[half:], 17, text_clr, font_name=body_font)

        else:  # "content"
            set_slide_background(slide, theme)
            add_accent_bar(slide, theme, 0.0, 0.0, 0.18, 7.5)
            if image_bytes:
                add_text(slide, 0.9, 0.6, 6.2, 1.0, title, 28, text_clr, bold=True, font_name=title_font)
                if subtitle:
                    add_text(slide, 0.9, 1.35, 6.2, 0.6, subtitle, 14, text_clr, font_name=body_font)
                add_bullets(slide, 0.9, 2.2, 6.2, 4.6, bullets, 16, text_clr, font_name=body_font)
                add_picture(slide, image_bytes, 7.6, 0.9, 4.8, 5.7)
            else:
                add_text(slide, 0.9, 0.6, 11.5, 1.0, title, 30, text_clr, bold=True, font_name=title_font)
                if subtitle:
                    add_text(slide, 0.9, 1.35, 11.5, 0.6, subtitle, 15, text_clr, font_name=body_font)
                add_bullets(slide, 0.9, 2.2, 11.5, 4.6, bullets, 18, text_clr, font_name=body_font)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# =====================================================================
# PREVIEW + DOWNLOADS
# =====================================================================
if st.session_state.slides:
    slides = st.session_state.slides
    slide_images = st.session_state.slide_images or {}
    html_deck = render_html_deck(slides, theme, slide_images)

    st.divider()
    st.subheader("👀 Preview")
    st.components.v1.html(html_deck, height=620, scrolling=False)

    st.divider()
    st.subheader("⬇️ Download your presentation")
    col1, col2 = st.columns(2)

    with col1:
        st.download_button(
            "Download HTML (dynamic deck)",
            data=html_deck.encode("utf-8"),
            file_name="presentation.html", mime="text/html", width="stretch",
        )

    with col2:
        if st.button("Prepare PPTX", width="stretch"):
            try:
                with st.spinner("Building your .pptx file..."):
                    st.session_state.pptx_bytes = build_pptx(slides, theme, slide_images)
            except Exception as e:
                st.error(f"PPTX generation failed: {e}")
        if st.session_state.pptx_bytes:
            st.download_button(
                "Download PPTX",
                data=st.session_state.pptx_bytes,
                file_name="presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                width="stretch",
            )

    st.caption(
        "The HTML deck and the PPTX file are both generated from the same slide "
        "content, so switching templates above updates both consistently."
    )
